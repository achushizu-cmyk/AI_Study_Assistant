"""
utils/file_parser.py
----------------------
Unified document/image text-extraction layer.

Every supported file type (PDF, DOCX, PPTX, TXT, JPG, PNG) is routed through
`extract_text()` which dispatches to the appropriate private extractor. This
is the ONLY module that should know about PyMuPDF / python-docx / python-pptx
internals - callers just get plain text back.
"""

import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

from utils.ocr_engine import OCREngine

_ocr_engine = None


def _get_ocr_engine() -> OCREngine:
    global _ocr_engine
    if _ocr_engine is None:
        _ocr_engine = OCREngine()
    return _ocr_engine


def _extract_from_pdf(file_path: str) -> str:
    """Extract text from a PDF. Falls back to OCR per-page for scanned PDFs."""
    text_chunks = []
    doc = fitz.open(file_path)
    try:
        for page in doc:
            page_text = page.get_text("text").strip()
            if page_text:
                text_chunks.append(page_text)
            else:
                # Likely a scanned/handwritten page - rasterize and OCR it.
                pix = page.get_pixmap(dpi=200)
                temp_image_path = file_path + f".page{page.number}.png"
                pix.save(temp_image_path)
                try:
                    ocr_text = _get_ocr_engine().extract_text(temp_image_path)
                    if ocr_text:
                        text_chunks.append(ocr_text)
                finally:
                    if os.path.exists(temp_image_path):
                        os.remove(temp_image_path)
    finally:
        doc.close()
    return "\n\n".join(text_chunks).strip()


def _extract_from_docx(file_path: str) -> str:
    document = Document(file_path)
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)
    return "\n".join(parts).strip()


def _extract_from_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)
    parts = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_lines = [f"--- Slide {slide_number} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        slide_lines.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_lines.append(row_text)
        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))
    return "\n\n".join(parts).strip()


def _extract_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def _extract_from_image(file_path: str) -> str:
    return _get_ocr_engine().extract_text(file_path)


_EXTRACTORS = {
    "pdf": _extract_from_pdf,
    "docx": _extract_from_docx,
    "pptx": _extract_from_pptx,
    "txt": _extract_from_txt,
    "jpg": _extract_from_image,
    "jpeg": _extract_from_image,
    "png": _extract_from_image,
}


def extract_text(file_path: str, file_extension: str) -> str:
    """
    Extract plain text from any supported file.

    Args:
        file_path: absolute path to the file on disk.
        file_extension: extension without the leading dot, e.g. "pdf".

    Returns:
        Extracted text, or an empty string if nothing could be extracted.

    Raises:
        ValueError: if the extension is not supported.
    """
    extension = file_extension.lower().lstrip(".")
    extractor = _EXTRACTORS.get(extension)
    if extractor is None:
        raise ValueError(f"Unsupported file type: .{extension}")
    return extractor(file_path)
